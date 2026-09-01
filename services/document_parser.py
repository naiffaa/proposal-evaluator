import csv
import hashlib
import html
import json
import re
import time
import unicodedata
import uuid
import xml.etree.ElementTree as ET
from pathlib import Path

import oci

from services.oci_auth import get_oci_auth_context


class DocumentParser:
    """
    Smart document parser.

    Fast local extraction:
    - PDF  -> PyMuPDF
    - DOCX -> python-docx
    - PPTX -> python-pptx
    - XLSX -> openpyxl
    - TXT / MD / CSV / JSON / XML / HTML / RTF -> local parsing

    OCI fallback:
    - Scanned / image-only PDFs
    - JPG / JPEG / PNG / TIF / TIFF

    Cache:
    - Successful extraction is cached by SHA-256.
    - Uploading the exact same file again skips extraction.
    """

    # ======================================================
    # Supported formats
    # ======================================================

    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".rtf",
    }

    OFFICE_EXTENSIONS = {
        ".docx",
        ".pptx",
        ".xlsx",
    }

    IMAGE_EXTENSIONS = {
        ".jpg",
        ".jpeg",
        ".png",
        ".tif",
        ".tiff",
    }

    SUPPORTED_EXTENSIONS = (
        TEXT_EXTENSIONS
        | OFFICE_EXTENSIONS
        | IMAGE_EXTENSIONS
        | {
            ".pdf",
        }
    )

    # ======================================================
    # Digital PDF detection
    # ======================================================

    MIN_PDF_TEXT_CHARS = 300

    MIN_AVG_CHARS_PER_PAGE = 80

    # ======================================================
    # Initialization
    # ======================================================

    def __init__(
        self,
        profile_name="DEFAULT",
        bucket_name="proposal-evaluator-documents",
        output_prefix="textExtraction",
        cache_directory=None,
    ):
        self.profile_name = (
            profile_name
        )

        self.bucket_name = (
            bucket_name
        )

        self.output_prefix = (
            output_prefix
        )

        # ==================================================
        # Cache directory
        # ==================================================

        project_root = (
            Path(__file__)
            .resolve()
            .parent
            .parent
        )

        if cache_directory:

            self.cache_directory = (
                Path(
                    cache_directory
                )
            )

        else:

            self.cache_directory = (
                project_root
                / "data"
                / "document_cache"
            )

        self.cache_directory.mkdir(
            parents=True,
            exist_ok=True,
        )

        # ==================================================
        # IMPORTANT:
        #
        # OCI clients are now initialized lazily.
        #
        # This means a normal digital PDF / Word / Excel
        # document does not connect to OCI at all.
        # ==================================================

        self.auth_mode = None

        self.signer = None

        self.region = None

        self.compartment_id = None

        self.config = None

        self.object_storage = None

        self.document_client = None

        self.namespace = None

        print()
        print(
            "--------------------------------"
        )

        print(
            "Smart Document Parser initialized"
        )

        print(
            "--------------------------------"
        )

        print(
            "Fast local extraction: enabled"
        )

        print(
            "OCI OCR fallback: enabled"
        )

        print(
            f"Cache: {self.cache_directory}"
        )

    # ======================================================
    # Validate file
    # ======================================================

    def _validate_file(
        self,
        file_path,
    ):
        file_path = (
            Path(
                file_path
            )
        )

        if not file_path.exists():

            raise FileNotFoundError(
                f"File does not exist: "
                f"{file_path}"
            )

        if not file_path.is_file():

            raise ValueError(
                f"Path is not a file: "
                f"{file_path}"
            )

        extension = (
            file_path
            .suffix
            .lower()
        )

        if (
            extension
            not in
            self.SUPPORTED_EXTENSIONS
        ):

            supported = (
                ", ".join(
                    sorted(
                        self.SUPPORTED_EXTENSIONS
                    )
                )
            )

            raise ValueError(
                "Unsupported document format.\n"
                f"Received: "
                f"{extension or 'no extension'}\n"
                f"Supported: {supported}"
            )

        return (
            file_path,
            extension,
        )

    # ======================================================
    # Text cleanup
    # ======================================================

    def _clean_text(
        self,
        text,
    ):
        if not isinstance(
            text,
            str,
        ):
            return ""

        # Arabic PDFs frequently extract as Unicode Arabic
        # Presentation Forms (U+FB50-U+FEFF), for example
        # "\ufedb\ufead\ufe8d\ufeb3\ufe94" instead of
        # "\u0643\u0631\u0627\u0633\u0629". Those code
        # points are visually identical but compare as
        # different strings, which silently breaks every
        # Arabic keyword match downstream (mandatory /
        # preferential labels, section detection, retrieval)
        # and wastes tokens in the model context.
        #
        # NFKC folds presentation forms back to standard
        # Arabic letters and leaves normal Latin/Arabic text
        # unchanged.
        if any(
            "\uFB50" <= character <= "\uFEFF"
            for character in text
        ):
            text = unicodedata.normalize(
                "NFKC",
                text,
            )

        text = (
            text
            .replace(
                "\x00",
                "",
            )
            .replace(
                "\r\n",
                "\n",
            )
            .replace(
                "\r",
                "\n",
            )
        )

        text = re.sub(
            r"[ \t]+\n",
            "\n",
            text,
        )

        text = re.sub(
            r"\n{4,}",
            "\n\n\n",
            text,
        )

        return (
            text.strip()
        )

    # ======================================================
    # SHA-256 cache
    # ======================================================

    def _calculate_file_hash(
        self,
        file_path,
    ):
        digest = (
            hashlib.sha256()
        )

        with open(
            file_path,
            "rb",
        ) as file_handle:

            while True:

                chunk = (
                    file_handle.read(
                        1024 * 1024
                    )
                )

                if not chunk:
                    break

                digest.update(
                    chunk
                )

        return (
            digest.hexdigest()
        )

    # Bump when text post-processing changes in a way that
    # makes previously cached extractions stale. Cached
    # entries written by an older version are ignored and
    # re-extracted rather than silently reused.
    #
    # v2: Arabic Presentation Forms are normalized to
    #     standard Arabic letters in _clean_text.
    TEXT_PIPELINE_VERSION = 2

    def _get_cache_path(
        self,
        file_hash,
    ):
        return (
            self.cache_directory
            / f"{file_hash}.json"
        )

    def _load_cached_result(
        self,
        file_path,
        file_hash,
    ):
        cache_path = (
            self._get_cache_path(
                file_hash
            )
        )

        if not cache_path.exists():
            return None

        try:

            cached = (
                json.loads(
                    cache_path.read_text(
                        encoding="utf-8",
                    )
                )
            )

        except Exception:

            return None

        cached_version = cached.get(
            "text_pipeline_version",
            1,
        )

        if (
            cached_version
            != self.TEXT_PIPELINE_VERSION
        ):
            print(
                "Cached extraction is stale "
                f"(v{cached_version} < "
                f"v{self.TEXT_PIPELINE_VERSION}). "
                "Re-extracting."
            )

            return None

        text = str(
            cached.get(
                "text",
                "",
            )
        ).strip()

        if not text:
            return None

        print()
        print(
            "--------------------------------"
        )

        print(
            "DOCUMENT CACHE HIT"
        )

        print(
            "--------------------------------"
        )

        print(
            f"File: "
            f"{Path(file_path).name}"
        )

        print(
            f"Method: "
            f"{cached.get('extraction_method')}"
        )

        print(
            f"Characters: "
            f"{len(text)}"
        )

        return {
            "file_name": (
                Path(
                    file_path
                ).name
            ),

            "file_hash": (
                file_hash
            ),

            "object_name": None,

            "processor_job_id": None,

            "result_prefix": None,

            "result_files": [],

            "text": (
                text
            ),

            "extraction_method": (
                cached.get(
                    "extraction_method",
                    "cache",
                )
            ),

            "page_count": (
                cached.get(
                    "page_count"
                )
            ),

            "cache_hit": True,

            "processing_time_seconds": 0.0,
        }

    def _save_cached_result(
        self,
        file_hash,
        result,
    ):
        cache_path = (
            self._get_cache_path(
                file_hash
            )
        )

        cache_data = {
            "file_hash": (
                file_hash
            ),

            "text_pipeline_version": (
                self.TEXT_PIPELINE_VERSION
            ),

            "text": (
                result.get(
                    "text",
                    "",
                )
            ),

            "extraction_method": (
                result.get(
                    "extraction_method",
                    "unknown",
                )
            ),

            "page_count": (
                result.get(
                    "page_count"
                )
            ),
        }

        cache_path.write_text(
            json.dumps(
                cache_data,
                indent=2,
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )

    # ======================================================
    # OCI lazy initialization
    # ======================================================

    def _ensure_oci_clients(
        self,
    ):
        if (
            self.object_storage
            is not None
            and
            self.document_client
            is not None
        ):
            return

        print()
        print(
            "--------------------------------"
        )

        print(
            "Initializing OCI OCR fallback"
        )

        print(
            "--------------------------------"
        )

        auth_context = (
            get_oci_auth_context()
        )

        self.auth_mode = (
            auth_context.mode
        )

        self.signer = (
            auth_context.signer
        )

        self.region = (
            auth_context.region
        )

        self.compartment_id = (
            auth_context.compartment_id
        )

        self.config = (
            auth_context.config
        )

        client_kwargs = {}

        if self.signer is not None:

            client_kwargs[
                "signer"
            ] = self.signer

        self.object_storage = (
            oci.object_storage
            .ObjectStorageClient(
                self.config,
                **client_kwargs,
            )
        )

        self.document_client = (
            oci.ai_document
            .AIServiceDocumentClient(
                self.config,
                **client_kwargs,
            )
        )

        namespace_response = (
            self.object_storage
            .get_namespace()
        )

        self.namespace = (
            namespace_response.data
        )

        print(
            f"Authentication: "
            f"{self.auth_mode}"
        )

        print(
            f"Region: "
            f"{self.region}"
        )

        print(
            f"Bucket: "
            f"{self.bucket_name}"
        )

        print(
            f"Namespace: "
            f"{self.namespace}"
        )

    # ======================================================
    # PDF - PyMuPDF
    # ======================================================

    def _extract_pdf_local(
        self,
        file_path,
    ):
        try:

            import pymupdf

        except ImportError:

            try:

                import fitz as pymupdf

            except ImportError as error:

                raise RuntimeError(
                    "PyMuPDF is not installed. "
                    "Run: pip install PyMuPDF"
                ) from error

        document = (
            pymupdf.open(
                str(
                    file_path
                )
            )
        )

        page_texts = []

        try:

            for page_index in range(
                len(
                    document
                )
            ):

                page = (
                    document[
                        page_index
                    ]
                )

                page_text = (
                    page.get_text(
                        "text"
                    )
                    or ""
                )

                page_text = (
                    self._clean_text(
                        page_text
                    )
                )

                if page_text:

                    page_texts.append(
                        (
                            f"[Page "
                            f"{page_index + 1}]\n"
                            f"{page_text}"
                        )
                    )

            final_text = (
                "\n\n"
                .join(
                    page_texts
                )
                .strip()
            )

            return {
                "text": (
                    final_text
                ),

                "page_count": (
                    len(
                        document
                    )
                ),
            }

        finally:

            document.close()

    # ======================================================
    # PDF quality check
    # ======================================================

    def _pdf_text_is_sufficient(
        self,
        text,
        page_count,
    ):
        text = (
            text
            or ""
        )

        compact_text = (
            re.sub(
                r"\s+",
                "",
                text,
            )
        )

        character_count = (
            len(
                compact_text
            )
        )

        if (
            character_count
            <
            self.MIN_PDF_TEXT_CHARS
        ):

            return False

        average_chars = (
            character_count
            /
            max(
                1,
                page_count,
            )
        )

        if (
            average_chars
            <
            self.MIN_AVG_CHARS_PER_PAGE
        ):

            return False

        return True

    # ======================================================
    # DOCX - Word
    # ======================================================

    def _extract_docx(
        self,
        file_path,
    ):
        try:

            from docx import Document

        except ImportError as error:

            raise RuntimeError(
                "python-docx is not installed. "
                "Run: pip install python-docx"
            ) from error

        document = (
            Document(
                str(
                    file_path
                )
            )
        )

        parts = []

        # --------------------------------------------------
        # Paragraphs
        # --------------------------------------------------

        for paragraph in (
            document.paragraphs
        ):

            text = (
                paragraph.text
                or ""
            ).strip()

            if text:

                parts.append(
                    text
                )

        # --------------------------------------------------
        # Tables
        # --------------------------------------------------

        for (
            table_index,
            table,
        ) in enumerate(
            document.tables,
            start=1,
        ):

            table_lines = [
                f"[Table {table_index}]"
            ]

            for row in (
                table.rows
            ):

                values = []

                for cell in (
                    row.cells
                ):

                    value = (
                        cell.text
                        or ""
                    ).strip()

                    values.append(
                        value
                    )

                if any(
                    values
                ):

                    table_lines.append(
                        " | ".join(
                            values
                        )
                    )

            if (
                len(
                    table_lines
                )
                >
                1
            ):

                parts.append(
                    "\n".join(
                        table_lines
                    )
                )

        final_text = (
            self._clean_text(
                "\n\n".join(
                    parts
                )
            )
        )

        return {
            "text": (
                final_text
            ),

            "page_count": None,
        }

    # ======================================================
    # PPTX - PowerPoint
    # ======================================================

    def _extract_pptx(
        self,
        file_path,
    ):
        try:

            from pptx import Presentation

        except ImportError as error:

            raise RuntimeError(
                "python-pptx is not installed. "
                "Run: pip install python-pptx"
            ) from error

        presentation = (
            Presentation(
                str(
                    file_path
                )
            )
        )

        slide_parts = []

        for (
            slide_index,
            slide,
        ) in enumerate(
            presentation.slides,
            start=1,
        ):

            parts = [
                f"[Slide {slide_index}]"
            ]

            for shape in (
                slide.shapes
            ):

                # ------------------------------------------
                # Normal text
                # ------------------------------------------

                shape_text = getattr(
                    shape,
                    "text",
                    None,
                )

                if (
                    isinstance(
                        shape_text,
                        str,
                    )
                    and
                    shape_text.strip()
                ):

                    parts.append(
                        shape_text.strip()
                    )

                # ------------------------------------------
                # Table
                # ------------------------------------------

                if getattr(
                    shape,
                    "has_table",
                    False,
                ):

                    for row in (
                        shape.table.rows
                    ):

                        values = [
                            (
                                cell.text
                                or ""
                            ).strip()

                            for cell
                            in row.cells
                        ]

                        if any(
                            values
                        ):

                            parts.append(
                                " | ".join(
                                    values
                                )
                            )

            if (
                len(
                    parts
                )
                >
                1
            ):

                slide_parts.append(
                    "\n".join(
                        parts
                    )
                )

        final_text = (
            self._clean_text(
                "\n\n".join(
                    slide_parts
                )
            )
        )

        return {
            "text": (
                final_text
            ),

            "page_count": (
                len(
                    presentation.slides
                )
            ),
        }

    # ======================================================
    # XLSX - Excel
    # ======================================================

    def _extract_xlsx(
        self,
        file_path,
    ):
        try:

            from openpyxl import (
                load_workbook,
            )

        except ImportError as error:

            raise RuntimeError(
                "openpyxl is not installed. "
                "Run: pip install openpyxl"
            ) from error

        workbook = (
            load_workbook(
                filename=str(
                    file_path
                ),
                read_only=True,
                data_only=True,
            )
        )

        sheet_parts = []

        try:

            for sheet in (
                workbook.worksheets
            ):

                lines = [
                    (
                        f"[Sheet: "
                        f"{sheet.title}]"
                    )
                ]

                for row in (
                    sheet.iter_rows(
                        values_only=True
                    )
                ):

                    values = []

                    for value in row:

                        if value is None:

                            values.append(
                                ""
                            )

                        else:

                            values.append(
                                str(
                                    value
                                ).strip()
                            )

                    if any(
                        values
                    ):

                        lines.append(
                            " | ".join(
                                values
                            )
                        )

                if (
                    len(
                        lines
                    )
                    >
                    1
                ):

                    sheet_parts.append(
                        "\n".join(
                            lines
                        )
                    )

        finally:

            workbook.close()

        final_text = (
            self._clean_text(
                "\n\n".join(
                    sheet_parts
                )
            )
        )

        return {
            "text": (
                final_text
            ),

            "page_count": None,
        }

    # ======================================================
    # Read plain text safely
    # ======================================================

    def _read_text_file(
        self,
        file_path,
    ):
        encodings = [
            "utf-8-sig",
            "utf-8",
            "utf-16",
            "cp1252",
        ]

        last_error = None

        for encoding in encodings:

            try:

                return (
                    file_path.read_text(
                        encoding=encoding
                    )
                )

            except UnicodeError as error:

                last_error = error

        raise RuntimeError(
            "Could not decode text file: "
            f"{file_path.name}"
        ) from last_error

    # ======================================================
    # TXT / MD
    # ======================================================

    def _extract_plain_text(
        self,
        file_path,
    ):
        text = (
            self._read_text_file(
                file_path
            )
        )

        return {
            "text": (
                self._clean_text(
                    text
                )
            ),

            "page_count": None,
        }

    # ======================================================
    # CSV
    # ======================================================

    def _extract_csv(
        self,
        file_path,
    ):
        raw_text = (
            self._read_text_file(
                file_path
            )
        )

        lines = []

        reader = (
            csv.reader(
                raw_text.splitlines()
            )
        )

        for row in reader:

            values = [
                str(
                    value
                ).strip()

                for value in row
            ]

            if any(
                values
            ):

                lines.append(
                    " | ".join(
                        values
                    )
                )

        return {
            "text": (
                self._clean_text(
                    "\n".join(
                        lines
                    )
                )
            ),

            "page_count": None,
        }

    # ======================================================
    # JSON
    # ======================================================

    def _extract_json(
        self,
        file_path,
    ):
        raw_text = (
            self._read_text_file(
                file_path
            )
        )

        try:

            data = (
                json.loads(
                    raw_text
                )
            )

            text = (
                json.dumps(
                    data,
                    indent=2,
                    ensure_ascii=False,
                )
            )

        except json.JSONDecodeError:

            text = (
                raw_text
            )

        return {
            "text": (
                self._clean_text(
                    text
                )
            ),

            "page_count": None,
        }

    # ======================================================
    # XML
    # ======================================================

    def _extract_xml(
        self,
        file_path,
    ):
        raw_text = (
            self._read_text_file(
                file_path
            )
        )

        try:

            root = (
                ET.fromstring(
                    raw_text
                )
            )

            text_parts = []

            for value in (
                root.itertext()
            ):

                value = (
                    value.strip()
                )

                if value:

                    text_parts.append(
                        value
                    )

            text = (
                "\n".join(
                    text_parts
                )
            )

        except ET.ParseError:

            text = (
                raw_text
            )

        return {
            "text": (
                self._clean_text(
                    text
                )
            ),

            "page_count": None,
        }

    # ======================================================
    # HTML
    # ======================================================

    def _extract_html(
        self,
        file_path,
    ):
        raw_text = (
            self._read_text_file(
                file_path
            )
        )

        try:

            from bs4 import (
                BeautifulSoup,
            )

            soup = (
                BeautifulSoup(
                    raw_text,
                    "html.parser",
                )
            )

            text = (
                soup.get_text(
                    "\n",
                    strip=True,
                )
            )

        except ImportError:

            text = (
                re.sub(
                    r"<[^>]+>",
                    " ",
                    raw_text,
                )
            )

            text = (
                html.unescape(
                    text
                )
            )

        return {
            "text": (
                self._clean_text(
                    text
                )
            ),

            "page_count": None,
        }

    # ======================================================
    # RTF
    # ======================================================

    def _extract_rtf(
        self,
        file_path,
    ):
        try:

            from striprtf.striprtf import (
                rtf_to_text,
            )

        except ImportError as error:

            raise RuntimeError(
                "striprtf is not installed. "
                "Run: pip install striprtf"
            ) from error

        raw_text = (
            self._read_text_file(
                file_path
            )
        )

        text = (
            rtf_to_text(
                raw_text
            )
        )

        return {
            "text": (
                self._clean_text(
                    text
                )
            ),

            "page_count": None,
        }

    # ======================================================
    # Upload file to OCI
    # ======================================================

    def upload_document(
        self,
        file_path,
    ):
        self._ensure_oci_clients()

        file_path = (
            Path(
                file_path
            )
        )

        extension = (
            file_path
            .suffix
            .lower()
        )

        if (
            extension
            not in
            (
                self.IMAGE_EXTENSIONS
                | {
                    ".pdf",
                }
            )
        ):

            raise ValueError(
                "OCI Document Understanding fallback "
                "supports PDF and image files only."
            )

        unique_id = (
            uuid.uuid4()
            .hex[:12]
        )

        safe_filename = (
            file_path.name
            .replace(
                "\\",
                "_",
            )
            .replace(
                "/",
                "_",
            )
        )

        object_name = (
            f"uploads/"
            f"{unique_id}_"
            f"{safe_filename}"
        )

        print()
        print(
            "--------------------------------"
        )

        print(
            "Uploading document to OCI"
        )

        print(
            "--------------------------------"
        )

        print(
            f"File: "
            f"{file_path.name}"
        )

        print(
            f"Object: "
            f"{object_name}"
        )

        try:

            upload_manager = (
                oci.object_storage
                .UploadManager(
                    self.object_storage
                )
            )

            upload_manager.upload_file(
                namespace_name=(
                    self.namespace
                ),
                bucket_name=(
                    self.bucket_name
                ),
                object_name=(
                    object_name
                ),
                file_path=str(
                    file_path
                ),
            )

        except oci.exceptions.ServiceError as error:

            raise RuntimeError(
                "Object Storage upload failed.\n"
                f"Status: {error.status}\n"
                f"Code: {error.code}\n"
                f"Message: {error.message}"
            ) from error

        print(
            "Upload successful."
        )

        return (
            object_name
        )

    # ======================================================
    # Start OCI Document Understanding
    # ======================================================

    def start_extraction_job(
        self,
        object_name,
    ):
        self._ensure_oci_clients()

        job_reference = (
            uuid.uuid4()
            .hex[:12]
        )

        result_prefix = (
            f"{self.output_prefix}/"
            f"{job_reference}"
        )

        input_location = (
            oci.ai_document.models
            .ObjectStorageLocations(
                object_locations=[
                    oci.ai_document.models
                    .ObjectLocation(
                        namespace_name=(
                            self.namespace
                        ),
                        bucket_name=(
                            self.bucket_name
                        ),
                        object_name=(
                            object_name
                        ),
                    )
                ]
            )
        )

        output_location = (
            oci.ai_document.models
            .OutputLocation(
                namespace_name=(
                    self.namespace
                ),
                bucket_name=(
                    self.bucket_name
                ),
                prefix=(
                    result_prefix
                ),
            )
        )

        text_feature = (
            oci.ai_document.models
            .DocumentTextExtractionFeature(
                generate_searchable_pdf=False,
                selection_mark_detection=False,
            )
        )

        processor_config = (
            oci.ai_document.models
            .GeneralProcessorConfig(
                features=[
                    text_feature
                ],
                document_type="OTHERS",
                is_zip_output_enabled=False,
            )
        )

        job_details = (
            oci.ai_document.models
            .CreateProcessorJobDetails(
                compartment_id=(
                    self.compartment_id
                ),
                input_location=(
                    input_location
                ),
                output_location=(
                    output_location
                ),
                processor_config=(
                    processor_config
                ),
                display_name=(
                    f"proposalEvaluator-"
                    f"{job_reference}"
                ),
            )
        )

        print()
        print(
            "--------------------------------"
        )

        print(
            "Starting OCI OCR fallback"
        )

        print(
            "--------------------------------"
        )

        try:

            response = (
                self.document_client
                .create_processor_job(
                    create_processor_job_details=(
                        job_details
                    )
                )
            )

        except oci.exceptions.ServiceError as error:

            raise RuntimeError(
                "Could not create "
                "Document Understanding job.\n"
                f"Status: {error.status}\n"
                f"Code: {error.code}\n"
                f"Message: {error.message}"
            ) from error

        job_id = (
            response.data.id
        )

        print(
            f"Processor Job ID: "
            f"{job_id}"
        )

        return (
            job_id,
            result_prefix,
        )

    # ======================================================
    # Wait for OCI job
    # ======================================================

    def wait_for_job(
        self,
        job_id,
        timeout_seconds=1200,
        polling_seconds=2,
    ):
        start_time = (
            time.time()
        )

        previous_state = None

        print()
        print(
            "--------------------------------"
        )

        print(
            "Waiting for OCI OCR"
        )

        print(
            "--------------------------------"
        )

        while True:

            try:

                response = (
                    self.document_client
                    .get_processor_job(
                        processor_job_id=(
                            job_id
                        )
                    )
                )

            except oci.exceptions.ServiceError as error:

                raise RuntimeError(
                    "Could not get "
                    "Document Understanding "
                    "job status.\n"
                    f"Status: {error.status}\n"
                    f"Code: {error.code}\n"
                    f"Message: {error.message}"
                ) from error

            job = (
                response.data
            )

            state = (
                job.lifecycle_state
            )

            if (
                state
                !=
                previous_state
            ):

                print(
                    f"Job State: "
                    f"{state}"
                )

                previous_state = (
                    state
                )

            if (
                state
                ==
                "SUCCEEDED"
            ):

                print(
                    "Analysis completed successfully."
                )

                return (
                    job
                )

            if state in {
                "FAILED",
                "CANCELED",
            }:

                details = getattr(
                    job,
                    "lifecycle_details",
                    None,
                )

                raise RuntimeError(
                    "Document Understanding "
                    "job failed.\n"
                    f"State: {state}\n"
                    f"Details: {details}"
                )

            elapsed = (
                time.time()
                -
                start_time
            )

            if (
                elapsed
                >
                timeout_seconds
            ):

                raise TimeoutError(
                    "Document analysis "
                    "took too long.\n"
                    f"Job ID: {job_id}\n"
                    f"Last State: {state}"
                )

            time.sleep(
                polling_seconds
            )

    # ======================================================
    # Find OCI result files
    # ======================================================

    def find_result_files(
        self,
        result_prefix,
    ):
        try:

            response = (
                oci.pagination
                .list_call_get_all_results(
                    self.object_storage
                    .list_objects,
                    namespace_name=(
                        self.namespace
                    ),
                    bucket_name=(
                        self.bucket_name
                    ),
                    prefix=(
                        result_prefix
                    ),
                )
            )

        except oci.exceptions.ServiceError as error:

            raise RuntimeError(
                "Could not list "
                "Document Understanding "
                "result files.\n"
                f"Status: {error.status}\n"
                f"Code: {error.code}\n"
                f"Message: {error.message}"
            ) from error

        result_objects = []

        for obj in (
            response.data.objects
        ):

            if (
                obj.name
                .lower()
                .endswith(
                    ".json"
                )
            ):

                result_objects.append(
                    obj.name
                )

        if not result_objects:

            raise RuntimeError(
                "OCI completed the job "
                "but no JSON result file "
                "was found."
            )

        result_objects.sort()

        return (
            result_objects
        )

    # ======================================================
    # Download OCI JSON
    # ======================================================

    def download_result_json(
        self,
        object_name,
    ):
        try:

            response = (
                self.object_storage
                .get_object(
                    namespace_name=(
                        self.namespace
                    ),
                    bucket_name=(
                        self.bucket_name
                    ),
                    object_name=(
                        object_name
                    ),
                )
            )

        except oci.exceptions.ServiceError as error:

            raise RuntimeError(
                "Could not download "
                "Document Understanding JSON.\n"
                f"Object: {object_name}\n"
                f"Status: {error.status}\n"
                f"Code: {error.code}\n"
                f"Message: {error.message}"
            ) from error

        content = (
            response.data.content
        )

        if isinstance(
            content,
            bytes,
        ):

            content = (
                content.decode(
                    "utf-8-sig"
                )
            )

        try:

            return (
                json.loads(
                    content
                )
            )

        except json.JSONDecodeError as error:

            raise RuntimeError(
                "Document Understanding "
                "output was not valid JSON."
            ) from error

    # ======================================================
    # Extract text from OCI result
    # ======================================================

    def extract_text(
        self,
        data,
    ):
        extracted_pages = []

        if (
            isinstance(
                data,
                dict,
            )
            and
            "analyzeDocumentResult"
            in data
        ):

            data = (
                data[
                    "analyzeDocumentResult"
                ]
            )

        if not isinstance(
            data,
            dict,
        ):

            return ""

        pages = (
            data.get(
                "pages",
                [],
            )
        )

        if not isinstance(
            pages,
            list,
        ):

            return ""

        for (
            page_index,
            page,
        ) in enumerate(
            pages,
            start=1,
        ):

            if not isinstance(
                page,
                dict,
            ):

                continue

            page_lines = []

            lines = (
                page.get(
                    "lines",
                    [],
                )
            )

            if isinstance(
                lines,
                list,
            ):

                for line in lines:

                    if not isinstance(
                        line,
                        dict,
                    ):

                        continue

                    text = (
                        line.get(
                            "text"
                        )
                    )

                    if (
                        isinstance(
                            text,
                            str,
                        )
                        and
                        text.strip()
                    ):

                        page_lines.append(
                            text.strip()
                        )

            # ----------------------------------------------
            # Fallback to words
            # ----------------------------------------------

            if not page_lines:

                words = (
                    page.get(
                        "words",
                        [],
                    )
                )

                page_words = []

                if isinstance(
                    words,
                    list,
                ):

                    for word in words:

                        if not isinstance(
                            word,
                            dict,
                        ):

                            continue

                        word_text = (
                            word.get(
                                "text"
                            )
                        )

                        if (
                            isinstance(
                                word_text,
                                str,
                            )
                            and
                            word_text.strip()
                        ):

                            page_words.append(
                                word_text.strip()
                            )

                if page_words:

                    page_lines.append(
                        " ".join(
                            page_words
                        )
                    )

            if page_lines:

                extracted_pages.append(
                    (
                        f"[Page "
                        f"{page_index}]\n"
                        +
                        "\n".join(
                            page_lines
                        )
                    )
                )

        return (
            self._clean_text(
                "\n\n".join(
                    extracted_pages
                )
            )
        )

    # ======================================================
    # Full OCI fallback workflow
    # ======================================================

    def _extract_with_oci(
        self,
        file_path,
    ):
        object_name = (
            self.upload_document(
                file_path
            )
        )

        (
            job_id,
            result_prefix,
        ) = (
            self.start_extraction_job(
                object_name
            )
        )

        self.wait_for_job(
            job_id
        )

        result_files = (
            self.find_result_files(
                result_prefix
            )
        )

        all_text = []

        for result_file in (
            result_files
        ):

            result_json = (
                self.download_result_json(
                    result_file
                )
            )

            text = (
                self.extract_text(
                    result_json
                )
            )

            if text:

                all_text.append(
                    text
                )

        final_text = (
            self._clean_text(
                "\n\n".join(
                    all_text
                )
            )
        )

        if not final_text:

            raise RuntimeError(
                "OCI analysis completed "
                "but no text was extracted."
            )

        return {
            "text": (
                final_text
            ),

            "page_count": None,

            "object_name": (
                object_name
            ),

            "processor_job_id": (
                job_id
            ),

            "result_prefix": (
                result_prefix
            ),

            "result_files": (
                result_files
            ),
        }

    # ======================================================
    # Main parsing workflow
    # ======================================================

    def parse_document(
        self,
        file_path,
    ):
        (
            file_path,
            extension,
        ) = (
            self._validate_file(
                file_path
            )
        )

        parse_started = (
            time.perf_counter()
        )

        # ==================================================
        # Hash + cache
        # ==================================================

        file_hash = (
            self._calculate_file_hash(
                file_path
            )
        )

        cached_result = (
            self._load_cached_result(
                file_path,
                file_hash,
            )
        )

        if (
            cached_result
            is not None
        ):

            return (
                cached_result
            )

        print()
        print(
            "=" * 60
        )

        print(
            "SMART DOCUMENT PROCESSING"
        )

        print(
            "=" * 60
        )

        print(
            f"File: "
            f"{file_path.name}"
        )

        print(
            f"Format: "
            f"{extension}"
        )

        result = None

        extraction_method = None

        # ==================================================
        # PDF
        # ==================================================

        if (
            extension
            ==
            ".pdf"
        ):

            print()
            print(
                "Trying fast PyMuPDF extraction..."
            )

            local_result = (
                self._extract_pdf_local(
                    file_path
                )
            )

            if (
                self._pdf_text_is_sufficient(
                    local_result.get(
                        "text",
                        "",
                    ),
                    local_result.get(
                        "page_count",
                        0,
                    ),
                )
            ):

                print(
                    "Digital PDF detected."
                )

                print(
                    "OCI Document Understanding skipped."
                )

                result = (
                    local_result
                )

                extraction_method = (
                    "pymupdf"
                )

            else:

                print(
                    "PDF appears scanned or "
                    "contains insufficient "
                    "extractable text."
                )

                print(
                    "Using OCI Document Understanding "
                    "as OCR fallback..."
                )

                result = (
                    self._extract_with_oci(
                        file_path
                    )
                )

                extraction_method = (
                    "oci_document_understanding"
                )

        # ==================================================
        # Images
        # ==================================================

        elif (
            extension
            in
            self.IMAGE_EXTENSIONS
        ):

            print(
                "Image document detected."
            )

            print(
                "Using OCI Document Understanding OCR..."
            )

            result = (
                self._extract_with_oci(
                    file_path
                )
            )

            extraction_method = (
                "oci_document_understanding"
            )

        # ==================================================
        # Word
        # ==================================================

        elif (
            extension
            ==
            ".docx"
        ):

            result = (
                self._extract_docx(
                    file_path
                )
            )

            extraction_method = (
                "python_docx"
            )

        # ==================================================
        # PowerPoint
        # ==================================================

        elif (
            extension
            ==
            ".pptx"
        ):

            result = (
                self._extract_pptx(
                    file_path
                )
            )

            extraction_method = (
                "python_pptx"
            )

        # ==================================================
        # Excel
        # ==================================================

        elif (
            extension
            ==
            ".xlsx"
        ):

            result = (
                self._extract_xlsx(
                    file_path
                )
            )

            extraction_method = (
                "openpyxl"
            )

        # ==================================================
        # CSV
        # ==================================================

        elif (
            extension
            ==
            ".csv"
        ):

            result = (
                self._extract_csv(
                    file_path
                )
            )

            extraction_method = (
                "local_csv"
            )

        # ==================================================
        # JSON
        # ==================================================

        elif (
            extension
            ==
            ".json"
        ):

            result = (
                self._extract_json(
                    file_path
                )
            )

            extraction_method = (
                "local_json"
            )

        # ==================================================
        # XML
        # ==================================================

        elif (
            extension
            ==
            ".xml"
        ):

            result = (
                self._extract_xml(
                    file_path
                )
            )

            extraction_method = (
                "local_xml"
            )

        # ==================================================
        # HTML
        # ==================================================

        elif (
            extension
            in
            {
                ".html",
                ".htm",
            }
        ):

            result = (
                self._extract_html(
                    file_path
                )
            )

            extraction_method = (
                "local_html"
            )

        # ==================================================
        # RTF
        # ==================================================

        elif (
            extension
            ==
            ".rtf"
        ):

            result = (
                self._extract_rtf(
                    file_path
                )
            )

            extraction_method = (
                "local_rtf"
            )

        # ==================================================
        # TXT / Markdown
        # ==================================================

        else:

            result = (
                self._extract_plain_text(
                    file_path
                )
            )

            extraction_method = (
                "local_text"
            )

        # ==================================================
        # Final text
        # ==================================================

        final_text = (
            self._clean_text(
                result.get(
                    "text",
                    "",
                )
            )
        )

        if not final_text:

            raise RuntimeError(
                "Document processing completed "
                "but no usable text was extracted."
            )

        elapsed = (
            time.perf_counter()
            -
            parse_started
        )

        final_result = {
            # =============================================
            # Existing keys kept for backward compatibility
            # =============================================

            "file_name": (
                file_path.name
            ),

            "object_name": (
                result.get(
                    "object_name"
                )
            ),

            "processor_job_id": (
                result.get(
                    "processor_job_id"
                )
            ),

            "result_prefix": (
                result.get(
                    "result_prefix"
                )
            ),

            "result_files": (
                result.get(
                    "result_files",
                    [],
                )
            ),

            "text": (
                final_text
            ),

            # =============================================
            # New diagnostic fields
            # =============================================

            "file_hash": (
                file_hash
            ),

            "extraction_method": (
                extraction_method
            ),

            "page_count": (
                result.get(
                    "page_count"
                )
            ),

            "cache_hit": False,

            "processing_time_seconds": (
                round(
                    elapsed,
                    3,
                )
            ),
        }

        # ==================================================
        # Cache successful result
        # ==================================================

        self._save_cached_result(
            file_hash,
            final_result,
        )

        print()
        print(
            "--------------------------------"
        )

        print(
            "Extraction complete"
        )

        print(
            "--------------------------------"
        )

        print(
            f"Method: "
            f"{extraction_method}"
        )

        print(
            f"Characters extracted: "
            f"{len(final_text)}"
        )

        print(
            f"Processing time: "
            f"{elapsed:.2f}s"
        )

        return (
            final_result
        )

    # ======================================================
    # Cleanup
    # ======================================================

    def close(
        self,
    ):
        for client in (
            self.document_client,
            self.object_storage,
        ):

            if client is None:

                continue

            close_method = getattr(
                client,
                "close",
                None,
            )

            if callable(
                close_method
            ):

                close_method()


# =========================================================
# Local Test
# =========================================================

if __name__ == "__main__":

    parser = (
        DocumentParser(
            bucket_name=(
                "proposal-evaluator-documents"
            ),
            output_prefix=(
                "textExtraction"
            ),
        )
    )

    test_file = (
        "Technical Proposal 01.pdf"
    )

    try:

        result = (
            parser.parse_document(
                test_file
            )
        )

        print()
        print(
            "=" * 60
        )

        print(
            "SMART DOCUMENT PARSER RESULT"
        )

        print(
            "=" * 60
        )

        print(
            f"File: "
            f"{result['file_name']}"
        )

        print(
            f"Method: "
            f"{result['extraction_method']}"
        )

        print(
            f"Cache Hit: "
            f"{result['cache_hit']}"
        )

        print(
            f"Processing Time: "
            f"{result['processing_time_seconds']}s"
        )

        print()
        print(
            "EXTRACTED TEXT PREVIEW"
        )

        print(
            "-" * 60
        )

        print(
            result[
                "text"
            ][:5000]
        )

    except Exception as error:

        print()
        print(
            "=" * 60
        )

        print(
            "ERROR"
        )

        print(
            "=" * 60
        )

        print(
            type(
                error
            ).__name__
        )

        print(
            str(
                error
            )
        )

        raise

    finally:

        parser.close()